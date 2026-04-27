"""
PPT Creator Service.
Handles AI-powered PowerPoint creation from uploaded PDF documents.
Reads PDF content, sends to AI for structuring into slides,
then builds a .pptx using varied layouts from the BSH template
with SmartArt-like diagrams and placeholder images.
"""
import io
import json
import math
import os
import re
import zipfile
from typing import List, Optional, Tuple

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app.services.common_service import (
    create_chat_history,
    call_brain_pure_llm_chat,
    upload_attachments,
    extract_pdf_text,
    sanitize_filename_for_prompt,
)

TEMPLATE_PATH = os.path.join(
    os.path.dirname(__file__), "..", "static", "docs", "pptTemplate.potx"
)

# ── Brand colours ──────────────────────────────────────────
BSH_ORANGE = RGBColor(0xFF, 0x5F, 0x00)
BSH_NAVY   = RGBColor(0x0F, 0x17, 0x2A)
BSH_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BSH_GREY   = RGBColor(0x64, 0x74, 0x8B)

# ── Layout name → index mapping (from the .potx template) ──
LAYOUT_MAP = {
    "title_slide":          0,
    "title_slide_with_image": 2,
    "chapter":              3,
    "content":              6,
    "smart_art":            13,
    "content_with_image":   11,
    "image_with_content":   12,
    "two_columns":          8,
    "three_columns":        9,
    "four_quadrants":       10,
    "full_image":           7,
    "title_only":           13,
    "end_slide":            15,
}

# ── SmartArt colour palette ───────────────────────────────
SMART_ART_COLORS = [
    RGBColor(0xFF, 0x5F, 0x00),  # BSH Orange
    RGBColor(0x22, 0x8B, 0xE6),  # Blue
    RGBColor(0x10, 0xB9, 0x81),  # Green
    RGBColor(0x8B, 0x5C, 0xF6),  # Purple
    RGBColor(0xF5, 0x9E, 0x0B),  # Amber
    RGBColor(0xEF, 0x44, 0x44),  # Red
    RGBColor(0x06, 0xB6, 0xD4),  # Cyan
    RGBColor(0x0F, 0x17, 0x2A),  # BSH Navy
]

ORANGE_THEME_HEX_COLORS = [
    "#FF5F00",  # Brand orange
    "#FF7A1A",  # Light orange
    "#E65100",  # Dark orange
    "#FFB380",  # Soft peach
    "#CC4C00",  # Burnt orange
    "#FFA64D",  # Amber-orange
]

# ── AI behaviour prompts ──────────────────────────────────

EXTRACT_BEHAVIOUR = (
    "You are a professional presentation designer. "
    "You receive text extracted from PDF documents and must structure it into "
    "a compelling PowerPoint presentation.\n\n"
    "AVAILABLE SLIDE LAYOUTS (use the exact layout name):\n"
    '- "title_slide": Opening slide with title & subtitle\n'
    '- "title_slide_with_image": Title slide with image placeholder (RARELY use - only when opening slide strongly benefits from a visual)\n'
    '- "chapter": Section divider / new chapter heading\n'
    '- "content": Standard slide with title and bullet points\n'
    '- "smart_art": Visual diagram instead of plain bullets (see SMART ART TYPES)\n'
    '- "content_with_image": Text on left + image placeholder on right\n'
    '- "image_with_content": Image placeholder on left + text on right\n'
    '- "two_columns": Two side-by-side content areas\n'
    '- "three_columns": Three side-by-side content areas\n'
    '- "four_quadrants": Four content areas in a 2x2 grid\n'
    '- "full_image": Title + full-width image placeholder\n'
    '- "title_only": Slide with only a large title\n'
    '- "end_slide": Closing / thank-you slide\n\n'
    "SMART ART TYPES (use with layout \"smart_art\"):\n"
    '- "process": Horizontal flow showing sequential steps (3-6 items)\n'
    '- "list_blocks": Vertical accent blocks for features/categories (3-6 items)\n'
    '- "pyramid": Layered pyramid for hierarchy/priority (3-5 items, top=highest)\n'
    '- "matrix": 2x2 grid for comparisons/categories (exactly 4 items)\n'
    '- "cycle": Circular arrangement for recurring processes (3-6 items)\n'
    '- "timeline": Horizontal timeline with events/milestones (3-6 items)\n'
    '- "venn": Overlapping circles showing relationships (2-3 items)\n'
    '- "funnel": Conversion funnel - wide at top, narrow at bottom (3-5 items)\n'
    '- "hierarchy": Org-chart tree — first item is root, rest are children (3-7 items)\n'
    '- "chevron_process": Chevron arrows with descriptions above/below. Items use pipe: "Step | Description" (3-6 items)\n'
    '- "timeline_detailed": Timeline with description boxes above/below. Items use pipe: "Milestone | Description" (3-6 items)\n'
    '- "agenda": Numbered agenda list with headlines. Items use pipe: "Headline | Subtext" (3-6 items)\n\n'
    "RULES:\n"
    "1. Start with a title_slide and end with an end_slide.\n"
    "2. Use chapter slides to separate major sections.\n (Optional)"
    "3. Use VARIED layouts — do NOT use the same layout for every slide. (Use what feels natural for the content)\n"
    "4. Use \"smart_art\" layout for AT LEAST 25-30% of content slides. (Make sure it fits the content, see types below)\n"
    "Choose the type that best fits:\n"
    "   - Sequential steps/workflows -> process or chevron_process (with descriptions)\n"
    "   - Features/specs/categories -> list_blocks\n"
    "   - Priority/hierarchy/levels -> pyramid\n"
    "   - 2x2 comparisons -> matrix\n"
    "   - Recurring/cyclical processes -> cycle\n"
    "   - Org structures/tree breakdowns -> hierarchy\n"
    "   - Events/milestones with details -> timeline_detailed\n"
    "   - Meeting agendas/action items -> agenda\n"
    "5. Use image placeholder slides (content_with_image, image_with_content, "
    "full_image) where a visual would enhance the message.\n"
    "6. Keep bullet points concise (max 6 per slide, max 15 words each).\n"
    "7. Smart art items should be short labels (2-5 words each).\n"
    "8. For two_columns / three_columns / four_quadrants, provide content "
    "in the 'columns' array.\n"
    "9. For smart_art, you can optionally specify custom 'colors' array with hex values "
    "(e.g., ['#FF5F00', '#228BE6']). If not specified, default palette is used.\n"
    "10. Never use the same 'layout' value for two or more consecutive slides. If you would "
    "repeat a 'content' slide back-to-back, switch to 'two_columns', 'smart_art', or a "
    "column variant instead. Visual variety is mandatory — reviewers notice monotony.\n\n"
    "Return your response as a valid JSON object with this EXACT structure:\n"
    "{\n"
    '  "title": "Presentation Title",\n'
    '  "subtitle": "Short subtitle",\n'
    '  "slides": [\n'
    "    {\n"
    '      "layout": "content",\n'
    '      "title": "Slide Title",\n'
    '      "subtitle": "Optional subtitle (title/chapter/end slides only)",\n'
    '      "bullets": ["Point 1", "Point 2"],\n'
    '      "smart_art": {"type": "process", "items": ["Step 1", "Step 2", "Step 3"], "colors": ["#FF5F00", "#228BE6"]},\n'
    '      "columns": [\n'
    '        {"heading": "Col 1", "bullets": ["item"]},\n'
    '        {"heading": "Col 2", "bullets": ["item"]}\n'
    "      ],\n"
    '      "image_description": "Brief description of what image should show",\n'
    '      "notes": "Speaker notes"\n'
    "    }\n"
    "  ]\n"
    "}\n\n"
    "FIELD USAGE BY LAYOUT:\n"
    "- title_slide / chapter / end_slide: title + subtitle\n"
    "- title_slide_with_image: title + subtitle + image_description\n"
    "- content: title + bullets\n"
    "- smart_art: title + smart_art object (type + items + optional colors)\n"
    "- content_with_image / image_with_content: title + bullets + image_description\n"
    "- two_columns / three_columns / four_quadrants: title + columns\n"
    "- full_image: title + image_description\n\n"
    "CRITICAL — USER INSTRUCTIONS OVERRIDE:\n"
    "If the user provides additional instructions (e.g., color preferences, slide count like 'one pager', "
    "specific themes, layout preferences, or any other customisation), you MUST follow them and they take "
    "HIGHEST PRIORITY over all other rules above. For example:\n"
    "- 'one pager' or 'single slide' = produce only 1-2 slides max with all key content condensed.\n"
    "- Color requests = apply the requested colors via smart_art.colors arrays and adjust content tone.\n"
    "- Slide count requests = strictly respect the requested number of slides.\n"
    "- Style/tone requests = adapt language, bullet length, and layout choices accordingly.\n\n"
    "Return ONLY the JSON, no markdown fences, no explanation."
)

REFINE_BEHAVIOUR = (
    "You are a PowerPoint content editor. "
    "The user wants changes to the presentation content. "
    "Apply the requested changes and return the FULL updated JSON "
    "with the same structure (title, subtitle, slides array with layout field). "
    "Keep using varied layouts including smart_art where appropriate. "
    "Available smart_art types: process, list_blocks, pyramid, matrix, cycle, timeline, venn, funnel, "
    "hierarchy, chevron_process, timeline_detailed, agenda. "
    "For types that support descriptions (chevron_process, timeline_detailed, agenda), "
    "items use pipe format: 'Label | Description'. "
    "For color changes, specify hex color codes in the smart_art.colors array. "
    "Return ONLY the JSON, no markdown fences, no explanation."
)


# ─── JSON parsing ────────────────────────────────────────

def _parse_json_response(text: str) -> dict:
    """Try to parse a JSON object from the AI response."""
    cleaned = re.sub(r"```(?:json)?", "", text).strip().rstrip("`").strip()

    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        pass

    match = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if match:
        try:
            return json.loads(match.group())
        except json.JSONDecodeError:
            pass

    return {}


# ─── Placeholder image generation ────────────────────────

def _create_placeholder_image(
    description: str = "Image Placeholder",
    width_px: int = 800,
    height_px: int = 600,
) -> io.BytesIO:
    """Create a clean light-grey placeholder image with centred description text."""
    img = Image.new("RGB", (width_px, height_px), (210, 210, 210))
    draw = ImageDraw.Draw(img)

    # Subtle darker-grey border
    draw.rectangle([0, 0, width_px - 1, height_px - 1], outline=(160, 160, 160), width=2)

    # Small camera/image icon hint — a thin inner rect to suggest a picture frame
    pad = width_px // 14
    draw.rectangle(
        [pad, pad, width_px - pad, height_px - pad],
        outline=(170, 170, 170), width=1,
    )

    # Description text — centred, dark grey
    try:
        font = ImageFont.truetype("arial.ttf", max(16, min(width_px, height_px) // 18))
    except (OSError, IOError):
        font = ImageFont.load_default()

    label = f"Photo: {description}"
    words = label.split()
    lines, current = [], ""
    for word in words:
        if len(current) + len(word) + 1 > 38:
            lines.append(current.strip())
            current = word + " "
        else:
            current += word + " "
    if current.strip():
        lines.append(current.strip())

    line_h = max(22, min(width_px, height_px) // 16)
    total_text_h = len(lines) * line_h
    text_y = (height_px - total_text_h) // 2
    for line in lines:
        try:
            lbbox = draw.textbbox((0, 0), line, font=font)
            lw = lbbox[2] - lbbox[0]
        except Exception:
            lw = len(line) * 8
        draw.text(((width_px - lw) // 2, text_y), line, fill=(90, 90, 90), font=font)
        text_y += line_h

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


# ─── SmartArt-like shape builders ────────────────────────

def _get_sa_color(idx: int, custom_colors: list = None) -> RGBColor:
    """Get a colour from custom list or default SmartArt palette."""
    if custom_colors and idx < len(custom_colors):
        hex_color = custom_colors[idx].lstrip('#')
        try:
            return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
        except (ValueError, IndexError):
            pass
    return SMART_ART_COLORS[idx % len(SMART_ART_COLORS)]


def _add_shape_with_text(
    slide, shape_type, left, top, width, height,
    text, fill_color, text_color=None, font_size=Pt(12), bold=True,
):
    """Add a shape with centred text."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    run = p.runs[0]
    run.font.size = font_size
    run.font.color.rgb = text_color or BSH_WHITE
    run.font.bold = bold
    return shape


def _draw_process_smart_art(slide, items, custom_colors=None):
    """Horizontal process flow — rounded boxes with arrows between them."""
    n = len(items)
    if n == 0:
        return

    area_left = Inches(0.6)
    area_top = Inches(1.8)
    area_w = Inches(10.8)
    area_h = Inches(4.0)

    arrow_w = Inches(0.28)
    gap = Inches(0.1)

    total_arrow_space = max(0, n - 1) * (arrow_w + gap * 2)
    box_w = min(int((area_w - total_arrow_space) / n), Inches(2.5))
    box_h = Inches(1.6)
    fs = Pt(12) if box_w > Inches(1.8) else Pt(11) if box_w > Inches(1.2) else Pt(10)

    total_w = n * box_w + max(0, n - 1) * (arrow_w + gap * 2)
    start_x = area_left + (area_w - total_w) // 2
    y = area_top + (area_h - box_h) // 2

    x = start_x
    for i, item in enumerate(items):
        color = _get_sa_color(i, custom_colors)
        _add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            int(x), int(y), int(box_w), int(box_h),
            item, color, BSH_WHITE, fs, True,
        )
        # Step number badge
        badge_r = Inches(0.22)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(x + Inches(0.08)), int(y + Inches(0.08)),
            int(badge_r * 2), int(badge_r * 2),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = BSH_WHITE
        badge.line.color.rgb = color
        badge.line.width = Pt(1)
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = str(i + 1)
        p_b.alignment = PP_ALIGN.CENTER
        for run in p_b.runs:
            run.font.size = Pt(8)
            run.font.color.rgb = color
            run.font.bold = True
        x += box_w
        if i < n - 1:
            x += gap
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                int(x), int(y + box_h // 2 - Inches(0.18)),
                int(arrow_w), int(Inches(0.36)),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BSH_ORANGE
            arrow.line.fill.background()
            x += arrow_w + gap


def _draw_list_blocks_smart_art(slide, items, custom_colors=None):
    """Vertical accent blocks with coloured left bar."""
    n = len(items)
    if n == 0:
        return

    area_left = Inches(1.5)
    area_top = Inches(1.8)
    area_w = Inches(9)
    area_h = Inches(4.5)

    gap = Inches(0.15)
    block_h = min(int((area_h - (n - 1) * gap) / n), Inches(1.1))
    accent_w = Inches(0.55)
    total_h = n * block_h + (n - 1) * gap
    y = area_top + (area_h - total_h) // 2
    fs_lb = Pt(16) if n <= 3 else Pt(14) if n <= 5 else Pt(12)
    badge_d = min(int(block_h * 0.55), int(Inches(0.45)))

    for i, item in enumerate(items):
        color = _get_sa_color(i, custom_colors)

        # Accent bar
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(area_left), int(y), int(accent_w), int(block_h),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        # Number badge centred on the accent bar
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(area_left + (accent_w - badge_d) // 2),
            int(y + (block_h - badge_d) // 2),
            int(badge_d), int(badge_d),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = BSH_WHITE
        badge.line.fill.background()
        tf_badge = badge.text_frame
        p_badge = tf_badge.paragraphs[0]
        p_badge.text = str(i + 1)
        p_badge.alignment = PP_ALIGN.CENTER
        for run in p_badge.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = color
            run.font.bold = True

        # Content block
        blk = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(area_left + accent_w + Inches(0.08)), int(y),
            int(area_w - accent_w - Inches(0.08)), int(block_h),
        )
        blk.fill.solid()
        blk.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        blk.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        blk.line.width = Pt(1)

        tf = blk.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = item
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = fs_lb
        run.font.color.rgb = BSH_NAVY
        run.font.bold = True

        y += block_h + gap


def _draw_pyramid_smart_art(slide, items, custom_colors=None):
    """Pyramid — widest at bottom, narrowest at top."""
    n = len(items)
    if n == 0:
        return

    area_left = Inches(1)
    area_top = Inches(1.8)
    area_w = Inches(10)
    area_h = Inches(4.5)

    gap = Inches(0.1)
    layer_h = min(int((area_h - (n - 1) * gap) / n), Inches(1))
    total_h = n * layer_h + (n - 1) * gap
    y = area_top + (area_h - total_h) // 2

    center_x = area_left + area_w // 2
    min_w = Inches(3)
    max_w = area_w

    for i, item in enumerate(items):
        frac = (n - 1 - i) / max(n - 1, 1)
        w = int(min_w + frac * (max_w - min_w))
        left = int(center_x - w // 2)
        _add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            left, int(y), w, int(layer_h),
            item, _get_sa_color(i, custom_colors), BSH_WHITE, Pt(12), True,
        )
        y += layer_h + gap


def _draw_matrix_smart_art(slide, items, custom_colors=None):
    """2x2 matrix grid."""
    padded = (items + ["", "", "", ""])[:4]

    area_left = Inches(1.5)
    area_top = Inches(1.8)
    area_w = Inches(9)
    area_h = Inches(4.2)

    gap = Inches(0.25)
    cell_w = int((area_w - gap) / 2)
    cell_h = int((area_h - gap) / 2)

    positions = [
        (area_left, area_top),
        (area_left + cell_w + gap, area_top),
        (area_left, area_top + cell_h + gap),
        (area_left + cell_w + gap, area_top + cell_h + gap),
    ]

    for i, (item, (x, y)) in enumerate(zip(padded, positions)):
        if not item:
            continue
        color = _get_sa_color(i, custom_colors)
        shape = _add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            int(x), int(y), cell_w, cell_h,
            item, color, BSH_WHITE, Pt(14), True,
        )
        shape.text_frame.margin_top = Inches(0.45)
        # Darker header strip at top of each cell
        darker = RGBColor(
            max(0, color[0] - 40),
            max(0, color[1] - 40),
            max(0, color[2] - 40),
        )
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(x), int(y), cell_w, int(Inches(0.38)),
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = darker
        hdr.line.fill.background()


def _draw_cycle_smart_art(slide, items, custom_colors=None):
    """Items in a circular arrangement representing a cycle."""
    n = len(items)
    if n == 0:
        return

    area_left = Inches(0.8)
    area_top = Inches(1.5)
    area_w = Inches(10.4)
    area_h = Inches(5)

    cx = area_left + area_w // 2
    cy = area_top + area_h // 2

    node_w = Inches(2)
    node_h = Inches(0.85)
    radius_x = (area_w - node_w) // 2 - Inches(0.3)
    radius_y = (area_h - node_h) // 2 - Inches(0.2)

    # Background circle outline
    cr = min(radius_x, radius_y) + Inches(0.15)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(cx - cr), int(cy - cr), int(cr * 2), int(cr * 2),
    )
    circle.fill.background()
    circle.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    circle.line.width = Pt(1.5)

    # Center anchor circle
    anchor_r = Inches(0.35)
    anchor = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(cx - anchor_r), int(cy - anchor_r),
        int(anchor_r * 2), int(anchor_r * 2),
    )
    anchor.fill.solid()
    anchor.fill.fore_color.rgb = BSH_ORANGE
    anchor.line.fill.background()

    for i, item in enumerate(items):
        angle = 2 * math.pi * i / n - math.pi / 2
        x = cx + radius_x * math.cos(angle) - node_w // 2
        y = cy + radius_y * math.sin(angle) - node_h // 2

        _add_shape_with_text(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            int(x), int(y), int(node_w), int(node_h),
            item, _get_sa_color(i, custom_colors), BSH_WHITE, Pt(11), True,
        )


def _draw_timeline_smart_art(slide, items, custom_colors=None):
    """Horizontal timeline with milestone markers."""
    n = len(items)
    if n == 0:
        return

    area_left = Inches(1)
    area_top = Inches(2.3)
    area_w = Inches(9)
    line_y = area_top + Inches(1)

    # Timeline line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(area_left), int(line_y - Inches(0.02)),
        int(area_w), int(Inches(0.04)),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BSH_GREY
    line.line.fill.background()

    gap = area_w / max(n - 1, 1) if n > 1 else 0
    marker_r = Inches(0.22)
    box_w_tl = Inches(1.6)
    box_h_tl = Inches(0.6)

    for i, item in enumerate(items):
        x = area_left + (gap * i if n > 1 else area_w // 2)
        color = _get_sa_color(i, custom_colors)

        # Marker circle
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(x - marker_r), int(line_y - marker_r),
            int(marker_r * 2), int(marker_r * 2),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.fill.background()

        # Label box with border, alternating above/below
        box_y = int(line_y - Inches(1.05)) if i % 2 == 0 else int(line_y + Inches(0.5))
        conn_top = int(box_y + box_h_tl) if i % 2 == 0 else int(line_y + marker_r)
        conn_bot = int(line_y - marker_r) if i % 2 == 0 else int(box_y)
        conn_h = abs(conn_bot - conn_top)
        if conn_h > 0:
            conn = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(x - Inches(0.015)), int(min(conn_top, conn_bot)),
                int(Inches(0.03)), conn_h,
            )
            conn.fill.solid()
            conn.fill.fore_color.rgb = color
            conn.line.fill.background()

        lbox = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(x - box_w_tl // 2), box_y,
            int(box_w_tl), int(box_h_tl),
        )
        lbox.fill.solid()
        lbox.fill.fore_color.rgb = BSH_WHITE
        lbox.line.color.rgb = color
        lbox.line.width = Pt(1.5)
        tf_l = lbox.text_frame
        tf_l.word_wrap = True
        tf_l.margin_left = Inches(0.08)
        tf_l.margin_top = Inches(0.06)
        p_l = tf_l.paragraphs[0]
        p_l.text = item
        p_l.alignment = PP_ALIGN.CENTER
        for run in p_l.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = BSH_NAVY
            run.font.bold = True


def _draw_venn_smart_art(slide, items, custom_colors=None):
    """Overlapping circles (2-3 circles) showing relationships."""
    n = min(len(items), 3)
    if n == 0:
        return

    area_left = Inches(2)
    area_top = Inches(2)
    circle_d = Inches(3)

    if n == 2:
        positions = [
            (area_left + Inches(1.5), area_top + Inches(1)),
            (area_left + Inches(3.5), area_top + Inches(1)),
        ]
    else:  # n == 3
        positions = [
            (area_left + Inches(2.5), area_top + Inches(0.5)),
            (area_left + Inches(1.5), area_top + Inches(2)),
            (area_left + Inches(3.5), area_top + Inches(2)),
        ]

    for i in range(n):
        x, y = positions[i]
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(x), int(y), int(circle_d), int(circle_d),
        )
        circle.fill.solid()
        color = _get_sa_color(i, custom_colors)
        circle.fill.fore_color.rgb = color
        circle.fill.transparency = 0.5
        circle.line.color.rgb = color
        circle.line.width = Pt(2)

        # Label
        tf = circle.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = items[i]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.color.rgb = BSH_NAVY
        p.runs[0].font.bold = True


def _draw_funnel_smart_art(slide, items, custom_colors=None):
    """Conversion funnel - wide at top, narrow at bottom."""
    n = len(items)
    if n == 0:
        return

    area_left = Inches(2)
    area_top = Inches(1.8)
    area_w = Inches(8)
    area_h = Inches(4.5)

    gap = Inches(0.1)
    stage_h = (area_h - (n - 1) * gap) / n

    y = area_top
    for i, item in enumerate(items):
        # Calculate width - narrow from top to bottom (wider bottom)
        frac = 1 - (i / max(n, 1)) * 0.35  # 100% to 65% width
        w = int(area_w * frac)
        x = int(area_left + (area_w - w) // 2)

        color = _get_sa_color(i, custom_colors)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, int(y), w, int(stage_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.6)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = item
        run = p.runs[0]
        run.font.size = Pt(13)
        run.font.color.rgb = BSH_WHITE
        run.font.bold = True

        # Number badge on the left edge of each stage
        badge_d = int(min(stage_h * 0.6, Inches(0.5)))
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(x + Inches(0.1)), int(y + (stage_h - badge_d) // 2),
            int(badge_d), int(badge_d),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = BSH_WHITE
        badge.line.fill.background()
        tf_badge = badge.text_frame
        p_badge = tf_badge.paragraphs[0]
        p_badge.text = str(i + 1)
        p_badge.alignment = PP_ALIGN.CENTER
        for run in p_badge.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = color
            run.font.bold = True

        y += stage_h + gap


def _draw_hierarchy_smart_art(slide, items, custom_colors=None):
    """Org-chart / hierarchy tree.

    Items convention: first item is root. Remaining items are children
    displayed in a grid below the root.  If >6 children, they wrap into
    two rows of up to 6 columns each (max 12 children).
    """
    n = len(items)
    if n == 0:
        return

    area_left = Inches(0.6)
    area_top = Inches(1.6)
    area_w = Inches(11)
    root_w = Inches(3.5)
    root_h = Inches(0.8)
    node_h = Inches(0.7)
    gap_x = Inches(0.2)
    gap_y = Inches(0.55)

    # Root node — centred at top
    root_x = area_left + (area_w - root_w) // 2
    root_y = area_top
    _add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        int(root_x), int(root_y), int(root_w), int(root_h),
        items[0], BSH_NAVY, BSH_WHITE, Pt(14), True,
    )

    children = items[1:]
    if not children:
        return

    # Split children into rows of up to 6
    max_per_row = 6
    rows = [children[i:i + max_per_row] for i in range(0, len(children), max_per_row)]

    row_top = root_y + root_h + gap_y
    root_cx = root_x + root_w // 2

    for row_idx, row_items in enumerate(rows):
        cols = len(row_items)
        node_w = min(int((area_w - (cols - 1) * gap_x) / cols), Inches(2.5))
        total_row_w = cols * node_w + (cols - 1) * gap_x
        start_x = area_left + (area_w - total_row_w) // 2
        current_row_y = row_top + row_idx * (node_h + gap_y + Inches(0.3))

        # Vertical connector from root to row
        if row_idx == 0:
            conn_x = int(root_cx)
            conn = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                conn_x - Inches(0.015), int(root_y + root_h),
                int(Inches(0.03)), int(current_row_y - (root_y + root_h)),
            )
            conn.fill.solid()
            conn.fill.fore_color.rgb = BSH_GREY
            conn.line.fill.background()

            # Horizontal bar connecting children
            first_cx = start_x + node_w // 2
            last_cx = start_x + (cols - 1) * (node_w + gap_x) + node_w // 2
            if cols > 1:
                hbar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    int(first_cx), int(current_row_y - Inches(0.15)),
                    int(last_cx - first_cx), int(Inches(0.03)),
                )
                hbar.fill.solid()
                hbar.fill.fore_color.rgb = BSH_GREY
                hbar.line.fill.background()

        x = start_x
        for ci, child in enumerate(row_items):
            color_idx = row_idx * max_per_row + ci + 1
            color = _get_sa_color(color_idx, custom_colors)

            # Vertical stub from horizontal bar to child
            child_cx = x + node_w // 2
            if row_idx == 0:
                stub = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    int(child_cx - Inches(0.015)), int(current_row_y - Inches(0.15)),
                    int(Inches(0.03)), int(Inches(0.15)),
                )
                stub.fill.solid()
                stub.fill.fore_color.rgb = BSH_GREY
                stub.line.fill.background()

            _add_shape_with_text(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                int(x), int(current_row_y), int(node_w), int(node_h),
                child, color, BSH_WHITE, Pt(11), True,
            )
            x += node_w + gap_x


def _draw_chevron_process_smart_art(slide, items, custom_colors=None):
    """Chevron arrow process with description boxes alternating above/below.

    Items are pairs: each item is either a plain string "Label" (rendered as
    chevron only) or "Label | Description" (chevron + description box).
    If no pipe, the whole string is the chevron label with no description.
    """
    n = len(items)
    if n == 0:
        return

    area_left = Inches(0.5)
    area_top = Inches(1.4)
    area_w = Inches(11.5)
    chevron_h = Inches(0.7)
    desc_box_h = Inches(1.1)
    desc_box_w = Inches(1.6)

    overlap = Inches(0.08)
    chev_w = min(int((area_w + (n - 1) * overlap) / n), Inches(2.5))
    total_w = n * chev_w - (n - 1) * overlap
    start_x = area_left + (area_w - total_w) // 2

    # Chevron strip sits in the vertical centre
    strip_y = area_top + desc_box_h + Inches(0.35)

    x = start_x
    for i, raw_item in enumerate(items):
        parts = raw_item.split("|", 1) if "|" in raw_item else [raw_item, ""]
        label = parts[0].strip()
        desc = parts[1].strip() if len(parts) > 1 else ""

        color = _get_sa_color(i, custom_colors)

        # Chevron shape
        _add_shape_with_text(
            slide, MSO_SHAPE.CHEVRON,
            int(x), int(strip_y), int(chev_w), int(chevron_h),
            label, color, BSH_WHITE, Pt(10), True,
        )

        # Description box (alternating above / below)
        if desc:
            chev_cx = x + chev_w // 2
            box_x = int(chev_cx - desc_box_w // 2)

            if i % 2 == 0:
                # Above
                box_y = int(strip_y - Inches(0.25) - desc_box_h)
                # Connector line
                conn = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    int(chev_cx - Inches(0.01)), int(box_y + desc_box_h),
                    int(Inches(0.02)), int(Inches(0.25)),
                )
            else:
                # Below
                box_y = int(strip_y + chevron_h + Inches(0.25))
                conn = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    int(chev_cx - Inches(0.01)), int(strip_y + chevron_h),
                    int(Inches(0.02)), int(Inches(0.25)),
                )

            conn.fill.solid()
            conn.fill.fore_color.rgb = BSH_GREY
            conn.line.fill.background()

            # Description text box
            txBox = slide.shapes.add_textbox(box_x, box_y, int(desc_box_w), int(desc_box_h))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_top = Inches(0.05)
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.05)
            p = tf.paragraphs[0]
            p.text = desc
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(9)
                run.font.color.rgb = BSH_NAVY

        x += chev_w - overlap


def _draw_timeline_detailed_smart_art(slide, items, custom_colors=None):
    """Detailed timeline — horizontal line with vertical connectors and
    description boxes alternating above and below each marker.

    Items: "Label | Description" or just "Label".
    """
    n = len(items)
    if n == 0:
        return

    area_left = Inches(0.5)
    area_w = Inches(11.5)
    line_y = Inches(3.6)
    box_w = Inches(1.5)
    box_h = Inches(1.2)
    connector_h = Inches(0.4)
    marker_r = Inches(0.1)

    # Horizontal line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(area_left), int(line_y - Inches(0.02)),
        int(area_w), int(Inches(0.04)),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BSH_ORANGE
    line.line.fill.background()

    spacing = area_w / max(n - 1, 1) if n > 1 else 0

    for i, raw_item in enumerate(items):
        parts = raw_item.split("|", 1) if "|" in raw_item else [raw_item, ""]
        label = parts[0].strip()
        desc = parts[1].strip() if len(parts) > 1 else ""

        cx = area_left + (spacing * i if n > 1 else area_w // 2)
        color = _get_sa_color(i, custom_colors)

        # Marker dot
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(cx - marker_r), int(line_y - marker_r),
            int(marker_r * 2), int(marker_r * 2),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.fill.background()

        # Vertical connector
        if i % 2 == 0:
            conn_top = line_y - marker_r - connector_h
            box_y = conn_top - box_h
        else:
            conn_top = line_y + marker_r
            box_y = conn_top + connector_h

        conn = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(cx - Inches(0.01)), int(conn_top),
            int(Inches(0.02)), int(connector_h),
        )
        conn.fill.solid()
        conn.fill.fore_color.rgb = color
        conn.line.fill.background()

        # Description box
        bx = int(cx - box_w // 2)
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bx, int(box_y), int(box_w), int(box_h),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        box.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        box.line.width = Pt(1)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.08)
        # Bold label paragraph
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = BSH_NAVY
            run.font.bold = True
        # Description paragraph
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.alignment = PP_ALIGN.LEFT
            for run in p2.runs:
                run.font.size = Pt(8)
                run.font.color.rgb = BSH_GREY


def _draw_agenda_smart_art(slide, items, custom_colors=None):
    """Numbered agenda list with headlines and optional sub-text.

    Items: "Headline | Sub-text" or just "Headline".
    Renders a numbered list on the left with separator lines, similar to
    the BSH Agenda template layout.
    """
    n = len(items)
    if n == 0:
        return

    area_left = Inches(0.8)
    area_top = Inches(1.6)
    area_w = Inches(10.5)
    row_h = Inches(0.85) if n <= 4 else Inches(0.75) if n == 5 else Inches(0.65)
    gap = Inches(0.08)
    num_w = Inches(0.6)
    y = area_top

    for i, raw_item in enumerate(items):
        parts = raw_item.split("|", 1) if "|" in raw_item else [raw_item, ""]
        headline = parts[0].strip()
        subtext = parts[1].strip() if len(parts) > 1 else ""
        color = _get_sa_color(i, custom_colors)

        # Top separator line
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(area_left), int(y),
            int(area_w), int(Inches(0.02)),
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        sep.line.fill.background()

        # Number
        num_box = slide.shapes.add_textbox(
            int(area_left), int(y + Inches(0.1)),
            int(num_w), int(row_h - Inches(0.1)),
        )
        tf_num = num_box.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = f"{i + 1}."
        p_num.alignment = PP_ALIGN.LEFT
        for run in p_num.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = color
            run.font.bold = True

        # Headline + subtext
        text_box = slide.shapes.add_textbox(
            int(area_left + num_w + Inches(0.15)), int(y + Inches(0.1)),
            int(area_w - num_w - Inches(0.15)), int(row_h - Inches(0.1)),
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = headline
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = BSH_NAVY
            run.font.bold = True

        if subtext:
            p2 = tf.add_paragraph()
            p2.text = subtext
            p2.alignment = PP_ALIGN.LEFT
            for run in p2.runs:
                run.font.size = Pt(12)
                run.font.color.rgb = BSH_GREY
                run.font.italic = True

        y += row_h + gap

    # Bottom separator line
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(area_left), int(y),
        int(area_w), int(Inches(0.02)),
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    sep.line.fill.background()


SMART_ART_DRAWERS = {
    "process": _draw_process_smart_art,
    "list_blocks": _draw_list_blocks_smart_art,
    "pyramid": _draw_pyramid_smart_art,
    "matrix": _draw_matrix_smart_art,
    "cycle": _draw_cycle_smart_art,
    "timeline": _draw_timeline_smart_art,
    "venn": _draw_venn_smart_art,
    "funnel": _draw_funnel_smart_art,
    "hierarchy": _draw_hierarchy_smart_art,
    "chevron_process": _draw_chevron_process_smart_art,
    "timeline_detailed": _draw_timeline_detailed_smart_art,
    "agenda": _draw_agenda_smart_art,
}


# ─── Template loader ─────────────────────────────────────

def _load_template() -> Presentation:
    """Load the .potx template, patching Content_Types for python-pptx compatibility."""
    if os.path.exists(TEMPLATE_PATH):
        raw = io.BytesIO()
        with zipfile.ZipFile(TEMPLATE_PATH, "r") as zin, zipfile.ZipFile(raw, "w") as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(
                        b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                        b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                    )
                zout.writestr(item, data)
        raw.seek(0)
        return Presentation(raw)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        return prs


# ─── Slide builders ──────────────────────────────────────

_EMOJI_RE = re.compile(
    r'[\U00010000-\U0010FFFF'     # Supplementary planes (most emoji)
    r'\u2600-\u26FF'              # Misc symbols
    r'\u2700-\u27BF'              # Dingbats
    r'\uFE00-\uFE0F'              # Variation selectors
    r'\u20D0-\u20FF]',            # Combining diacritical marks for symbols
    re.UNICODE,
)

_ADML_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_RPR_TAGS = frozenset([
    f'{{{_ADML_NS}}}rPr',
    f'{{{_ADML_NS}}}defRPr',
    f'{{{_ADML_NS}}}endParaRPr',
])


def _force_no_strike(tf):
    """Force strike='noStrike' on every run/paragraph property element in a text frame.

    Some BSH template placeholders have strikethrough baked into the placeholder XML at
    the paragraph or body level. python-pptx's run.font.strike only sets it on the run
    element; iterating the raw XML ensures no level in the inheritance chain has strike.
    """
    for elem in tf._txBody.iter():
        if elem.tag in _RPR_TAGS:
            elem.set('strike', 'noStrike')


def _clean_text(s: str) -> str:
    """Strip emoji and markdown formatting markers from text before inserting into slides."""
    if not s:
        return s
    s = re.sub(r'~~(.*?)~~', r'\1', s)          # ~~strikethrough~~ -> text
    s = re.sub(r'\*\*(.*?)\*\*', r'\1', s)     # **bold** -> text
    s = re.sub(r'\*(.*?)\*', r'\1', s)          # *italic* -> text
    s = _EMOJI_RE.sub('', s)
    return s.strip()


def _set_placeholder_text(slide, ph_idx: int, text: str):
    """Safely set text on a placeholder by index."""
    try:
        slide.placeholders[ph_idx].text = _clean_text(text)
    except (KeyError, IndexError):
        pass


def _fill_bullets(placeholder, bullets: list[str]):
    """Fill a content placeholder with bullet points, auto-scaling font by count."""
    try:
        n = len(bullets)
        fs = Pt(18) if n <= 3 else Pt(16) if n <= 5 else Pt(14)
        tf = placeholder.text_frame
        tf.clear()
        for j, bullet in enumerate(bullets):
            text = _clean_text(bullet)
            if j == 0:
                tf.paragraphs[0].text = text
                _style_para(tf.paragraphs[0], fs, BSH_NAVY)
            else:
                p = tf.add_paragraph()
                p.text = text
                _style_para(p, fs, BSH_NAVY)
        _force_no_strike(tf)
    except Exception:
        pass


def _insert_placeholder_image(slide, ph_idx: int, description: str):
    """Insert a placeholder image into a picture placeholder."""
    try:
        pic_ph = slide.placeholders[ph_idx]
        img_buf = _create_placeholder_image(description)
        pic_ph.insert_picture(img_buf)
    except (KeyError, IndexError, Exception):
        pass


def _add_textbox(slide, text, left, top, width, height, font_size, bold, color):
    """Add a simple textbox to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    _style_para(p, font_size, color, bold)


def _style_para(para, size, color, bold=False):
    """Apply font styling to a paragraph."""
    for run in para.runs:
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.strike = False


def _add_title_accent(slide):
    """Draw a short orange underline beneath the slide title."""
    try:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.87), Inches(1.43),
            Inches(2.5), Inches(0.045),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = BSH_ORANGE
        bar.line.fill.background()
    except Exception:
        pass


def _add_chapter_accent(slide):
    """Draw an orange vertical accent bar on the left of a chapter slide."""
    try:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.3), Inches(1.8),
            Inches(0.15), Inches(3.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = BSH_ORANGE
        bar.line.fill.background()
    except Exception:
        pass


def _build_title_slide(prs, slide_data, layout, is_chapter: bool = False):
    """Build a title / chapter / end slide."""
    slide = prs.slides.add_slide(layout)
    _set_placeholder_text(slide, 0, slide_data.get("title", ""))
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        _set_placeholder_text(slide, 1, subtitle)
    return slide


def _build_content_slide(prs, slide_data, layout):
    """Build a standard content slide with title + bullets."""
    slide = prs.slides.add_slide(layout)
    _set_placeholder_text(slide, 0, slide_data.get("title", ""))
    bullets = slide_data.get("bullets", [])
    if bullets:
        try:
            _fill_bullets(slide.placeholders[1], bullets)
        except (KeyError, IndexError):
            y = Inches(1.6)
            for b in bullets:
                _add_textbox(slide, f"• {b}", Inches(1), y, Inches(11), Inches(0.5), Pt(16), False, BSH_NAVY)
                y += Inches(0.55)
    _add_title_accent(slide)
    return slide


def _build_content_with_image_slide(prs, slide_data, layout):
    """Build text+image slide (layout 11: text left, picture right)."""
    slide = prs.slides.add_slide(layout)
    _set_placeholder_text(slide, 0, slide_data.get("title", ""))
    bullets = slide_data.get("bullets", [])
    if bullets:
        try:
            _fill_bullets(slide.placeholders[1], bullets)
        except (KeyError, IndexError):
            pass
    desc = slide_data.get("image_description", "Relevant visual")
    _insert_placeholder_image(slide, 2, desc)
    _add_title_accent(slide)
    return slide


def _build_image_with_content_slide(prs, slide_data, layout):
    """Build image+text slide (layout 12: picture left, text right)."""
    slide = prs.slides.add_slide(layout)
    _set_placeholder_text(slide, 0, slide_data.get("title", ""))
    bullets = slide_data.get("bullets", [])
    if bullets:
        try:
            _fill_bullets(slide.placeholders[2], bullets)
        except (KeyError, IndexError):
            pass
    desc = slide_data.get("image_description", "Relevant visual")
    _insert_placeholder_image(slide, 1, desc)
    _add_title_accent(slide)
    return slide


def _build_full_image_slide(prs, slide_data, layout):
    """Build title + full image slide (layout 7)."""
    slide = prs.slides.add_slide(layout)
    _set_placeholder_text(slide, 0, slide_data.get("title", ""))
    desc = slide_data.get("image_description", "Full-width visual")
    _insert_placeholder_image(slide, 1, desc)
    _add_title_accent(slide)
    return slide


def _build_multi_column_slide(prs, slide_data, layout, num_cols: int):
    """Build a multi-column slide (2, 3, or 4 columns)."""
    slide = prs.slides.add_slide(layout)
    _set_placeholder_text(slide, 0, slide_data.get("title", ""))

    columns = slide_data.get("columns", [])

    for col_idx in range(num_cols):
        ph_idx = col_idx + 1
        if col_idx < len(columns):
            col_data = columns[col_idx]
            heading = col_data.get("heading", "")
            col_bullets = col_data.get("bullets", [])
            all_text = []
            if heading:
                all_text.append(_clean_text(heading))
            all_text.extend([_clean_text(b) for b in col_bullets])
            if all_text:
                try:
                    tf = slide.placeholders[ph_idx].text_frame
                    tf.clear()
                    for j, txt in enumerate(all_text):
                        if j == 0:
                            tf.paragraphs[0].text = txt
                            _style_para(
                                tf.paragraphs[0],
                                Pt(16) if j == 0 and heading else Pt(14),
                                BSH_ORANGE if bool(heading) else BSH_NAVY,
                                j == 0 and bool(heading),
                            )
                        else:
                            p = tf.add_paragraph()
                            p.text = txt
                            _style_para(p, Pt(14), BSH_NAVY)
                except (KeyError, IndexError):
                    pass

    return slide


def _build_title_only_slide(prs, slide_data, layout):
    """Build a title-only slide (layout 13)."""
    slide = prs.slides.add_slide(layout)
    _set_placeholder_text(slide, 0, slide_data.get("title", ""))
    return slide


def _build_smart_art_slide(prs, slide_data, layout, force_orange_theme: bool = False):
    """Build a slide with SmartArt-like diagram shapes."""
    slide = prs.slides.add_slide(layout)
    _set_placeholder_text(slide, 0, slide_data.get("title", ""))

    sa = slide_data.get("smart_art", {})
    sa_type = sa.get("type", "list_blocks")
    sa_items = [_clean_text(item) for item in sa.get("items", [])]
    sa_colors = sa.get("colors", None)

    # If orange theme is forced and AI didn't provide explicit colors, apply orange palette.
    if force_orange_theme and not sa_colors:
        sa_colors = ORANGE_THEME_HEX_COLORS

    drawer = SMART_ART_DRAWERS.get(sa_type)
    if drawer and sa_items:
        drawer(slide, sa_items, sa_colors)

    return slide    


# ─── Main PPT builder ────────────────────────────────────

def generate_pptx_file(
    content: dict,
    username: str = "Unknown User",
    force_orange_theme: bool = False,
) -> io.BytesIO:
    """
    Build a .pptx from structured AI content using varied template layouts.
    Inserts SmartArt diagrams and placeholder images where indicated.
    Adds footer with date and username to each slide.
    """
    from datetime import datetime
    
    prs = _load_template()
    
    # Remove any default slides that come with the template
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    slides_data = content.get("slides", [])
    current_date = datetime.now().strftime("%B %d, %Y")  # e.g., "March 01, 2026"
    short_date = datetime.now().strftime("%Y-%m-%d")      # e.g., "2026-03-21"
    pres_title = content.get("title", "")

    # Update master slide footer: date and presentation title
    for master in prs.slide_masters:
        for shape in master.shapes:
            if not shape.has_text_frame:
                continue
            if shape.name == "DateOnSlides":
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = short_date
            elif shape.name == "FooterOnSlides":
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = pres_title

    for slide_data in slides_data:
        layout_name = slide_data.get("layout", "content")
        layout_idx = LAYOUT_MAP.get(layout_name, LAYOUT_MAP["content"])

        if layout_idx >= len(prs.slide_layouts):
            layout_idx = LAYOUT_MAP["content"]

        layout = prs.slide_layouts[layout_idx]

        if layout_name in ("title_slide", "chapter", "end_slide"):
            slide = _build_title_slide(prs, slide_data, layout, is_chapter=(layout_name == "chapter"))

        elif layout_name == "title_slide_with_image":
            slide = _build_title_slide(prs, slide_data, layout)
            # Add image placeholder if description provided
            desc = slide_data.get("image_description", "Opening visual")
            _insert_placeholder_image(slide, 2, desc)

        elif layout_name == "content":
            slide = _build_content_slide(prs, slide_data, layout)

        elif layout_name == "smart_art":
            slide = _build_smart_art_slide(
                prs,
                slide_data,
                layout,
                force_orange_theme=force_orange_theme,
            )

        elif layout_name == "content_with_image":
            slide = _build_content_with_image_slide(prs, slide_data, layout)

        elif layout_name == "image_with_content":
            slide = _build_image_with_content_slide(prs, slide_data, layout)

        elif layout_name == "full_image":
            slide = _build_full_image_slide(prs, slide_data, layout)

        elif layout_name == "two_columns":
            slide = _build_multi_column_slide(prs, slide_data, layout, 2)

        elif layout_name == "three_columns":
            slide = _build_multi_column_slide(prs, slide_data, layout, 3)

        elif layout_name == "four_quadrants":
            slide = _build_multi_column_slide(prs, slide_data, layout, 4)

        elif layout_name == "title_only":
            slide = _build_title_only_slide(prs, slide_data, layout)

        else:
            slide = _build_content_slide(prs, slide_data, layout)

        notes = slide_data.get("notes", "")
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass
        
        # Add footer with date and username (bottom-left) - only on title_slide and end_slide
        if layout_name in ("title_slide", "title_slide_with_image", "end_slide"):
            try:
                footer_text = f"{current_date}\n{username}"
                footer_box = slide.shapes.add_textbox(
                    Inches(0.3), 
                    prs.slide_height - Inches(0.6),
                    Inches(2.5), 
                    Inches(0.4)
                )
                tf = footer_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = footer_text
                for run in p.runs:
                    run.font.size = Pt(8)
                    run.font.color.rgb = BSH_WHITE  # White color
            except Exception:
                pass  # Silently fail if footer can't be added

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ─── Public API ─────────────────────────────────────────────

async def extract_pdf_content(
    pdf_bytes_list: list[tuple[str, bytes]],
    user_instructions: str = "",
    username: str = "Unknown User",
    image_files: list = None,
    force_orange_theme: bool = False,
) -> dict:
    """
    Extract text from uploaded PDFs and/or process uploaded images via AI vision,
    then structure content into presentation slides.
    Returns dict with 'content' (parsed JSON), 'chatHistoryId', or 'error'.
    """
    brain_id = os.getenv("DSCP_BRAIN_ID")
    if not brain_id:
        return {"error": True, "message": "API Not Active", "detail": "DSCP_BRAIN_ID is not configured."}

    all_text_parts: list[str] = []
    errors: list[str] = []
    for fname, fbytes in pdf_bytes_list:
        text, err = extract_pdf_text(fbytes)
        if err:
            errors.append(f"{sanitize_filename_for_prompt(fname)}: {err}")
        elif text:
            all_text_parts.append(f"=== File: {sanitize_filename_for_prompt(fname)} ===\n{text}")

    has_text = bool(all_text_parts)
    has_images = bool(image_files)

    if not has_text and not has_images:
        detail = "\n".join(errors) if errors else "No readable content could be extracted from the uploaded files."
        return {
            "error": True,
            "message": "Extraction Failed",
            "detail": detail,
        }

    combined_text = "\n\n".join(all_text_parts) if has_text else ""

    base_task = (
        "Create a professional PowerPoint presentation from this content.\n"
        "Use varied slide layouts and SmartArt diagrams to make it visually engaging.\n"
        "Include image placeholder slides where visuals would enhance the message."
    )

    if has_text:
        body = f"Below is text content extracted from uploaded documents.\n{base_task}\n\n{combined_text}"
    else:
        body = f"Analyze the uploaded image(s) and {base_task}"

    if force_orange_theme:
        body = (
            "VISUAL THEME REQUIREMENT: Use an orange-first visual style. "
            "Prefer orange shades for smart_art.colors and supportive warm accents where appropriate.\n\n"
            f"{body}"
        )

    if user_instructions:
        prompt = (
            f"\u26a0\ufe0f USER INSTRUCTIONS (HIGHEST PRIORITY \u2014 follow these strictly):\n{user_instructions}\n\n"
            f"---\n\n{body}"
        )
    else:
        prompt = body

    chat_result = await create_chat_history(brain_id)
    if chat_result.get("error"):
        return chat_result

    chat_history_id = chat_result.get("chatHistoryId")

    attachment_ids = []
    if has_images:
        upload_result = await upload_attachments(brain_id, image_files)
        if upload_result.get("error"):
            return upload_result
        attachment_ids = upload_result.get("attachmentIds", [])

    response = await call_brain_pure_llm_chat(
        brain_id,
        prompt,
        chat_history_id=chat_history_id,
        attachment_ids=attachment_ids or None,
        custom_behaviour=EXTRACT_BEHAVIOUR,
    )

    raw = response.get("result", "")
    parsed = _parse_json_response(raw)

    if not parsed or "slides" not in parsed:
        return {
            "error": True,
            "message": "Could not parse AI response",
            "detail": "The AI did not return valid slide JSON.",
            "raw": raw,
        }

    return {
        "content": parsed,
        "chatHistoryId": response.get("chatHistoryId", chat_history_id),
        "username": username,
    }


async def refine_ppt_content(
    chat_history_id: str,
    message: str,
    current_content: dict = None,
    force_orange_theme: bool = False,
) -> dict:
    """Continue the conversation to refine the presentation content."""
    brain_id = os.getenv("DSCP_BRAIN_ID")
    if not brain_id:
        return {"error": True, "message": "API Not Active", "detail": "DSCP_BRAIN_ID is not configured."}

    theme_prefix = ""
    if force_orange_theme:
        theme_prefix = (
            "VISUAL THEME REQUIREMENT: Keep an orange-first visual style and use orange shades "
            "in smart_art.colors where relevant.\n\n"
        )

    prompt = f"{theme_prefix}{message}" if theme_prefix else message
    if current_content:
        prompt = (
            f"{theme_prefix}"
            f"Here is the current presentation JSON:\n{json.dumps(current_content, indent=2)}\n\n"
            f"User request: {message}\n\n"
            "Apply the changes and return the FULL updated JSON with the same structure "
            "(title, subtitle, slides array with layout field)."
        )

    response = await call_brain_pure_llm_chat(
        brain_id,
        prompt,
        chat_history_id=chat_history_id,
        custom_behaviour=REFINE_BEHAVIOUR,
    )

    if response.get("error"):
        return response

    raw = response.get("result", "")
    parsed = _parse_json_response(raw)

    if not parsed or "slides" not in parsed:
        return {
            "result": raw,
            "chatHistoryId": response.get("chatHistoryId", chat_history_id),
        }

    return {
        "content": parsed,
        "chatHistoryId": response.get("chatHistoryId", chat_history_id),
    }
